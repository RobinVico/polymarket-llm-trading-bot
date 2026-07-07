#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
策略总览 PPT 生成器 (企业风格 / 可打印, 16:9)。

内容 = v7.4.4 现行策略快照 (2026-07-06 从代码核对):
  monitor.py 出场常量 / auto_reeval.py 触发+护栏+模型 / sizing.py 金额公式 /
  scanner.py FILTERS / 实时持仓快照举例。
⚠️ 阈值是生成时手工抄自代码的快照 —— 以后版本变了要先核对更新本文件再重新生成。

依赖 python-pptx (项目 .venv 里没有, 用任意装了它的 python 跑):
    python3 -m venv /tmp/pptenv && /tmp/pptenv/bin/pip install python-pptx
    /tmp/pptenv/bin/python3 scripts/gen_strategy_ppt.py
输出: 项目根目录 Polymarket_Bot_策略总览_v7.4.4.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

VERSION = "7.4.4"
DATE = "2026-07-06"
OUT = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                   f"Polymarket_Bot_策略总览_v{VERSION}.pptx")

# ---------- 版式常量 ----------
PAGE_W, PAGE_H = Inches(13.333), Inches(7.5)
ML = Inches(0.55)            # 左右边距
CW = Inches(12.233)          # 内容宽

NAVY = '1B2A4A'; NAVY2 = '24395F'; ACCENT = '2E74B5'; GOLD = 'C9A227'
INK = '222B36'; GRAY = '66707F'; FAINT = '8A93A3'
WHITE = 'FFFFFF'; ZEBRA = 'F3F5F9'
GREEN = '1E7B4F'; GREEN_BG = 'E3F3EA'
RED = 'B3362B'; RED_BG = 'FBE7E4'
AMBER = '8F6212'; AMBER_BG = 'FFF2D9'
BLUE_BG = 'E8F0FA'; CHIP_BG = 'F5F7FA'
KICKER_C = '9FB6D9'

_page = {'n': 0}


def R(h):
    return RGBColor.from_string(h)


def _fmt(run, size=13, bold=False, color=INK, italic=False, font='Helvetica Neue'):
    f = run.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.name = font
    f.color.rgb = R(color)
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        latin = rPr.find(qn('a:latin'))
        if latin is not None:
            latin.addnext(ea)
        else:
            rPr.append(ea)
    ea.set('typeface', 'PingFang SC')


def txt(slide, x, y, w, h, paras):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = para.get('align', PP_ALIGN.LEFT)
        p.space_after = Pt(para.get('space_after', 0))
        p.space_before = Pt(para.get('space_before', 0))
        p.line_spacing = para.get('line', 1.12)
        for rr in para['runs']:
            r = p.add_run()
            r.text = rr[0]
            _fmt(r, **(rr[1] if len(rr) > 1 else {}))
    return tb


def rect(slide, x, y, w, h, fill, rounded=False, radius=0.10, line_color=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = R(fill)
    if line_color:
        shp.line.color.rgb = R(line_color); shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    return shp


def bullets(slide, x, y, w, items, size=12.3, gap=6, line=1.16, marker='•', marker_color=ACCENT):
    paras = []
    for it in items:
        pre = None; mk = marker; mkc = marker_color; sz = size
        if isinstance(it, str):
            body = it
        elif isinstance(it, tuple):
            pre, body = it
        else:
            pre = it.get('b'); body = it['t']
            mk = it.get('m', marker); mkc = it.get('mc', marker_color); sz = it.get('size', size)
        runs = [(mk + '  ', {'size': sz, 'bold': True, 'color': mkc})]
        if pre:
            runs.append((pre, {'size': sz, 'bold': True, 'color': INK}))
        runs.append((body, {'size': sz, 'color': INK}))
        paras.append({'runs': runs, 'space_after': gap, 'line': line})
    return txt(slide, x, y, w, Inches(0.4), paras)


def table(slide, x, y, w, col_ws, data, fs=11.5, header_fs=None, row_h=0.5, header_h=0.42,
          col_align=None, header_fill=NAVY, first_col_bold=True, cell_over=None):
    """data: 第一行是表头; 单元格 = str (支持 \n, 第2行起小号灰) 或 {'lines':[(text,opts),...]}.
    cell_over: {(ri,ci): {run opts}} 覆盖正文格式 (如标红)."""
    rows = len(data); cols = len(data[0])
    gf = slide.shapes.add_table(rows, cols, x, y, w, Inches(header_h + row_h * (rows - 1)))
    t = gf.table
    t.first_row = False; t.horz_banding = False
    for ci, cw_in in enumerate(col_ws):
        t.columns[ci].width = Inches(cw_in)
    t.rows[0].height = Inches(header_h)
    for ri in range(1, rows):
        t.rows[ri].height = Inches(row_h)
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            c = t.cell(ri, ci)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = Inches(0.1); c.margin_right = Inches(0.07)
            c.margin_top = Inches(0.03); c.margin_bottom = Inches(0.03)
            c.fill.solid()
            c.fill.fore_color.rgb = R(header_fill if ri == 0 else (WHITE if ri % 2 == 1 else ZEBRA))
            tf = c.text_frame; tf.word_wrap = True
            if isinstance(val, dict):
                lines = val['lines']
            else:
                lines = [(s, None) for s in val.split('\n')]
            for li, (text, opts) in enumerate(lines):
                p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                p.alignment = (col_align[ci] if col_align else PP_ALIGN.LEFT)
                p.line_spacing = 1.05
                r = p.add_run(); r.text = text
                if opts is not None:
                    _fmt(r, **opts)
                elif ri == 0:
                    _fmt(r, size=header_fs or fs, bold=True, color=WHITE)
                else:
                    o = {'size': fs if li == 0 else fs - 1.5,
                         'bold': first_col_bold and ci == 0 and li == 0,
                         'color': INK if li == 0 else GRAY}
                    if cell_over and (ri, ci) in cell_over:
                        o.update(cell_over[(ri, ci)])
                    _fmt(r, **o)
    return gf


def callout(slide, x, y, w, h, title, body, bar=RED, bg=RED_BG, body_size=11.3, title_size=12.3):
    rect(slide, x, y, w, h, bg, rounded=True, radius=0.08)
    rect(slide, x, y + Inches(0.09), Inches(0.075), h - Inches(0.18), bar)
    paras = [{'runs': [(title, {'size': title_size, 'bold': True, 'color': bar})], 'space_after': 3, 'line': 1.1}]
    for ln in body.split('\n'):
        paras.append({'runs': [(ln, {'size': body_size, 'color': INK})], 'line': 1.18, 'space_after': 1})
    txt(slide, x + Inches(0.24), y + Inches(0.1), w - Inches(0.42), h - Inches(0.2), paras)


def chips(slide, x, y, w, items, h=1.0, gap=0.24, num_size=19, cap_size=10.5,
          fill=CHIP_BG, num_color=NAVY, cap_color=GRAY, line_color='DDE3EC'):
    n = len(items)
    bw = (w - Inches(gap) * (n - 1)) / n
    for i, (num, cap) in enumerate(items):
        bx = x + i * (bw + Inches(gap))
        rect(slide, bx, y, bw, h, fill, rounded=True, radius=0.12, line_color=line_color)
        txt(slide, bx + Inches(0.08), y + Inches(0.16), bw - Inches(0.16), Inches(0.4), [
            {'runs': [(num, {'size': num_size, 'bold': True, 'color': num_color})], 'align': PP_ALIGN.CENTER}])
        txt(slide, bx + Inches(0.08), y + h - Inches(0.42), bw - Inches(0.16), Inches(0.3), [
            {'runs': [(cap, {'size': cap_size, 'color': cap_color})], 'align': PP_ALIGN.CENTER}])


def flow(slide, x, y, w, h, steps, fill=NAVY2, head_size=12, sub_size=9.5,
         head_color=WHITE, sub_color='C9D6EC', arrow_zone=0.34):
    n = len(steps)
    bw = (w - Inches(arrow_zone) * (n - 1)) / n
    for i, (head, sub) in enumerate(steps):
        bx = x + i * (bw + Inches(arrow_zone))
        rect(slide, bx, y, bw, h, fill, rounded=True, radius=0.1)
        paras = [{'runs': [(head, {'size': head_size, 'bold': True, 'color': head_color})],
                  'align': PP_ALIGN.CENTER, 'space_after': 3, 'line': 1.05}]
        for ln in sub.split('\n'):
            paras.append({'runs': [(ln, {'size': sub_size, 'color': sub_color})],
                          'align': PP_ALIGN.CENTER, 'line': 1.1})
        txt(slide, bx + Inches(0.06), y + Inches(0.13), bw - Inches(0.12), h - Inches(0.26), paras)
        if i < n - 1:
            txt(slide, bx + bw, y + h / 2 - Inches(0.16), Inches(arrow_zone), Inches(0.32), [
                {'runs': [('➜', {'size': 13, 'bold': True, 'color': ACCENT})], 'align': PP_ALIGN.CENTER}])


def new_slide(prs, kicker, title):
    _page['n'] += 1
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, PAGE_W, Inches(1.02), NAVY)
    rect(s, ML, Inches(0.30), Inches(0.085), Inches(0.46), GOLD)
    txt(s, ML + Inches(0.22), Inches(0.12), CW, Inches(0.28), [
        {'runs': [(kicker, {'size': 10.5, 'bold': True, 'color': KICKER_C})]}])
    txt(s, ML + Inches(0.22), Inches(0.36), CW, Inches(0.5), [
        {'runs': [(title, {'size': 21.5, 'bold': True, 'color': WHITE})]}])
    rect(s, ML, Inches(7.05), CW, Inches(0.014), 'D5DAE3')
    txt(s, ML, Inches(7.12), Inches(8), Inches(0.25), [
        {'runs': [(f'Polymarket 交易机器人 · 策略总览 v{VERSION} · {DATE}', {'size': 8.5, 'color': FAINT})]}])
    txt(s, PAGE_W - Inches(3.55), Inches(7.12), Inches(3), Inches(0.25), [
        {'runs': [(f'第 {_page["n"]:02d} 页 · 内部资料', {'size': 8.5, 'color': FAINT})],
         'align': PP_ALIGN.RIGHT}])
    return s


# ==================================================================
prs = Presentation()
prs.slide_width = PAGE_W
prs.slide_height = PAGE_H
prs.core_properties.title = f'Polymarket Bot 策略总览 v{VERSION}'
prs.core_properties.author = 'Polymarket Bot'

# ---------- 1. 封面 ----------
s = prs.slides.add_slide(prs.slide_layouts[6]); _page['n'] += 1
rect(s, 0, 0, PAGE_W, PAGE_H, NAVY)
rect(s, 0, Inches(4.92), PAGE_W, Inches(0.02), GOLD)
txt(s, ML, Inches(0.42), Inches(6), Inches(0.3), [
    {'runs': [('内部资料 · 打印版', {'size': 10.5, 'color': KICKER_C})]}])
vchip = rect(s, PAGE_W - Inches(1.9), Inches(0.38), Inches(1.35), Inches(0.44), GOLD, rounded=True, radius=0.5)
txt(s, PAGE_W - Inches(1.9), Inches(0.45), Inches(1.35), Inches(0.3), [
    {'runs': [(f'v{VERSION}', {'size': 13, 'bold': True, 'color': NAVY})], 'align': PP_ALIGN.CENTER}])
txt(s, ML, Inches(1.85), CW, Inches(0.9), [
    {'runs': [('Polymarket 交易机器人', {'size': 40, 'bold': True, 'color': WHITE})]}])
txt(s, ML, Inches(2.78), CW, Inches(0.6), [
    {'runs': [('现行策略总览', {'size': 25, 'bold': True, 'color': GOLD})]}])
txt(s, ML, Inches(3.42), CW, Inches(0.4), [
    {'runs': [('止盈止损 · 自动重评 · 仓位金额 · 找仓规则 · 资金风控', {'size': 14.5, 'color': KICKER_C})]}])
txt(s, ML, Inches(4.02), CW, Inches(0.35), [
    {'runs': [(f'版本 v{VERSION} · 生成于 2026 年 7 月 6 日 · 封面数据为生成时刻实时快照',
               {'size': 11, 'color': '8DA3C6'})]}])
chips(s, ML, Inches(5.35), CW, [
    ('$93.44', '总资产'),
    ('$65.71', '现金'),
    ('5 笔 · $27.73', '持仓数 · 持仓市值'),
    ('30 秒', '盯盘心跳'),
], h=Inches(1.05), fill=NAVY2, num_color=WHITE, cap_color=KICKER_C, line_color=None)

# ---------- 2. 目录 ----------
s = new_slide(prs, 'CONTENTS', '目录')
toc = [
    ('01', '系统总览 — 一页看懂整个流程', 3),
    ('02', '止盈 — 什么时候锁定利润', 4),
    ('03', '止损 — 什么时候认输离场', 5),
    ('04', '其他出场 & 每 30 秒的决策顺序', 6),
    ('05', '自动重评 ① 什么时候触发', 7),
    ('06', '自动重评 ② AI 怎么决定、谁来执行', 8),
    ('07', '自动重评 ③ 用哪个 AI', 9),
    ('08', '新仓金额 — 一笔下多少钱', 10),
    ('09', '找新仓 — 从扫描到录入', 11),
    ('10', '入场门槛 — edge 要多大才推荐', 12),
    ('11', '钱的总规矩 — 风控红线一页', 13),
    ('12', '测试仓 /paper — 不花钱的模拟盘', 14),
    ('13', '常改参数速查表', 15),
    ('14', '已知风险与注意事项', 16),
]
for half, x0 in ((toc[:7], ML), (toc[7:], Inches(7.0))):
    paras = []
    for num, title_, pg in half:
        paras.append({'runs': [
            (num, {'size': 13.5, 'bold': True, 'color': GOLD, 'font': 'Menlo'}),
            ('   ' + title_, {'size': 13.5, 'color': INK}),
            (f'   ·  P{pg}', {'size': 11, 'color': FAINT}),
        ], 'space_after': 15})
    txt(s, x0, Inches(1.55), Inches(5.9), Inches(5), paras)
callout(s, Inches(7.0), Inches(5.55), Inches(5.78), Inches(1.28),
        '怎么用这份文件',
        '平时只需要翻两页: P13 (钱的红线) 和 P15 (参数速查)。\n出了问题先看 P16 (已知风险)。想改规则时, 对应页脚注了参数在哪改。',
        bar=ACCENT, bg=BLUE_BG)

# ---------- 3. 系统总览 ----------
s = new_slide(prs, '01 · 系统总览', '一页看懂: 从找市场到卖出')
flow(s, ML, Inches(1.28), CW, Inches(1.18), [
    ('① 扫市场', 'bot 按流动性/时间窗\n过滤 39 个主题'),
    ('② Claude 分析', '给方向 · 胜率 q\n信心 · 止损档 · 主题簇'),
    ('③ 算金额', '公式给建议\n每仓 $1–15'),
    ('④ 人工下单', '你去 Polymarket 手动买\nbot 没有自动买入'),
    ('⑤ bot 盯盘', '每 30 秒评估一次\n该卖时自动卖'),
])
bullets(s, ML, Inches(2.78), CW, [
    ('分工: ', 'Claude 出主意 · bot 盯盘和执行卖出 · 买入永远你亲手下单 — bot 代码里没有自动买入路径'),
    ('卖出可以自动: ', '止盈到线 / 跌破 $0.05 地板 / 你离线时重评结论=卖 — 规则详见 P4–P9'),
    ('节奏: ', '30 秒一次决策心跳 · 30 分钟一次资产快照+自动备份 · 1 小时一次结算检查'),
    ('每个仓位录入时定好三件事, bot 全靠它们盯盘: ', '胜率 q · 止损档 (三选一, 见下) · 主题簇 (防同类扎堆)'),
    ('几块屏: ', '主页 (桌面操作) · /panel 副屏控制台 · /m 手机只读 · /paper 测试仓 · /history 往期复盘'),
], size=12.6, gap=9)
callout(s, ML, Inches(5.55), CW, Inches(1.22),
        '三种止损档 (入场时 Claude 分类, 决定这个仓后面怎么止盈止损)',
        '收敛型 convergent = 真相会自己收敛的题 (票房/数据/汇率/比分) · 混合型 hybrid = 民调+政治混合 (选举类)\n事件型 event_driven = 政治/外交/谈判 — 价格天天震荡但震荡≠真相变化, 所以规则对它最宽容',
        bar=NAVY, bg=CHIP_BG, body_size=11.6)

# ---------- 4. 止盈 ----------
s = new_slide(prs, '02 · 出场策略', '止盈 — 什么时候锁定利润')
table(s, ML, Inches(1.26), CW, [3.5, 4.0, 4.73], [
    ['仓位类型', '触发条件 (按"真能卖到的价")', '动作'],
    ['事件型 · 翻倍先到', '卖价 ≥ 2×成本 (且早于 $0.92)', '全部卖出锁翻倍 (低价入场时先到)'],
    ['事件型 · 到 $0.92', '卖价 ≥ $0.92 (还没翻倍)', '卖一半留一半; 后半跌破 $0.78 再全卖'],
    ['收敛型 · 距结算 ≤3 天', '卖价 ≥ $0.88', '全部卖出 (提前落袋, 留出滑点)'],
    ['收敛型 >3 天 / 混合型', '卖价 ≥ $0.90 或 浮盈 ≥ +100%', '全部卖出'],
], fs=12, row_h=0.46)
bullets(s, ML, Inches(3.95), CW, [
    ('为什么看"卖价"? ', '触发用盘口买一价 best_bid (你真能成交的价), 不用参考价 — 防"显示 $0.90 实际只能卖 $0.60"的假止盈'),
    ('事件型这套怎么跑? ', '论点没破常一路涨到底。低价捡的、翻倍(≥2×成本)先到 → 全卖落袋; 高价的到 $0.92 → 卖一半留一半博结算, 但留的半仓从 0.92 跌 15%(<$0.78)就把后半也卖了锁利润, 不让它坐过山车吐回去'),
    ('可以手动关: ', '每仓有"关止盈"开关, 关了只跳过止盈, 止损照常跑'),
], size=12.2, gap=8)
callout(s, ML, Inches(5.62), CW, Inches(1.18),
        '拿现在的仓举例',
        'SpaceX 星舰 No (事件型) $0.72 → 若先翻倍到 2×成本就全卖锁翻倍; 否则涨到 $0.92 卖一半留一半 (后半跌破 $0.78 再清)。\n霍尔木兹海峡恢复 No (混合型) $0.855 → 涨到 $0.90 会自动全部卖出。',
        bar=GREEN, bg=GREEN_BG)

# ---------- 5. 止损 ----------
s = new_slide(prs, '03 · 出场策略', '止损 — 什么时候认输离场')
table(s, ML, Inches(1.26), CW, [2.9, 4.6, 4.73], [
    ['仓位类型', '止损方式', '触发线'],
    ['收敛型 convergent', '移动止损: 盯"持有期最高价", 从峰值回撤', '回撤 ≥20% (距结算 ≤3 天收紧到 12%)\n且连续 6 次心跳 ≈3 分钟确认, 防一抖卖飞'],
    ['混合型 hybrid', '移动止损: 盯"持有期最高价", 从峰值回撤 (v7.4.4)', '回撤 ≥35% 且连续 6 次心跳确认 (跟收敛型同形式)'],
    ['事件型 event_driven', '很松的 %止损 + 地板 (震荡 ≠ 真相变化)', '亏 ≥60% (v7.4.3 新增) 或 跌破 $0.05 地板'],
    ['未分类 (少见)', '默认当"混合型"处理 (不再有 -25% 老仓档)', '回撤 ≥35% (同混合型移动止损)'],
], fs=12, row_h=0.52)
callout(s, ML, Inches(4.14), CW, Inches(0.95),
        '关键: 砸穿止损线 ≠ 马上卖!',
        '先进「⏸ 等重评」状态, 让 AI 联网重查后再决定 (P7–P8)。只有两种情况不问 AI 直接卖: ① 跌破 $0.05 地板 ② 重评功能没配置。',
        bar=RED, bg=RED_BG)
bullets(s, ML, Inches(5.32), CW, [
    ('两种止损口径: ', '收敛型 + 混合型盯"持有期最高价"回撤 (移动止损); 事件型盯"成本价"(实时加权均价, 加仓自动摊平) 亏 60%。成本价永远用实时均价, 不用第一次入场价'),
    ('事件型止损为什么这么松: ', '政治/外交题天天震荡, 按普通 % 卖会卖错 (回测 3 笔 2 笔卖完就反弹); 但也不能只靠地板一路跌到底, 所以 v7.4.3 给个很松的 -60% 兜底'),
    ('活例子: ', 'Trump 与 MBS 通话 Yes (事件型) 现在 −55% — 没跌破地板就不硬卖; 今天已自动重评, 按新胜率 40% 继续持有 (仍比现价高 17 个点)'),
], size=12.2, gap=8)

# ---------- 6. 其他出场 + 优先级 ----------
s = new_slide(prs, '04 · 出场策略', '其他出场 & 每 30 秒的决策顺序')
txt(s, ML, Inches(1.22), Inches(6.2), Inches(0.35), [
    {'runs': [('每 30 秒对每个仓从上往下过一遍, 第一个命中的执行:', {'size': 12.8, 'bold': True, 'color': NAVY})]}])
bullets(s, ML, Inches(1.66), Inches(6.15), [
    {'m': '①', 'mc': GOLD, 'b': '没录入信息', 't': ' → 状态 NO_META, 提醒补录 (不动仓)'},
    {'m': '②', 'mc': GOLD, 'b': '止盈', 't': ' (P4) — 含事件型 翻倍先到全卖 / 0.92 卖一半'},
    {'m': '③', 'mc': GOLD, 'b': '止损 / 移动止损', 't': ' (P5) — 多数进「⏸ 等重评」'},
    {'m': '④', 'mc': GOLD, 'b': '时间止损: ', 't': '距结算 ≤2 天且价格离入场价没挪动 (<5 个点) → 自动全卖 — 快结算了还在原地 = 论点没兑现, 别让钱陪跑到最后一刻'},
    {'m': '⑤', 'mc': GOLD, 'b': '状态灯', 't': ' (只亮灯提示, 不自动卖) — 见右表'},
], size=12.4, gap=10)
table(s, Inches(7.05), Inches(1.32), Inches(5.73), [1.85, 2.05, 1.83], [
    ['状态灯', '胜率 q − 现价', '意思'],
    ['HOLD', '> +2 点', '继续拿'],
    ['MARGINAL', '−3 ～ +2 点', '边缘, 别加仓'],
    ['AT_TARGET', '< −3 点 (没重评过)', '已到目标价, 考虑落袋'],
    ['SOFT_NEGATIVE', '< −3 点 (重评过)', 'AI 也偏空, 人工留意'],
], fs=10.8, row_h=0.5, header_h=0.38)
bullets(s, Inches(7.05), Inches(4.1), Inches(5.73), [
    '状态灯每 30 秒刷新, 显示在主页和副屏每仓的「决策状态」一行',
    ('活例子: ', 'Wesley Bell 仓现在就是 MARGINAL — q 44% vs 现价 43% 只差 1 个点, 拿着但别加'),
], size=11.4, gap=7)
callout(s, ML, Inches(5.62), CW, Inches(1.18),
        '记住一条: 触发价口径不对称 (防假信号的关键设计)',
        '止盈看 best_bid (真能卖到的价) — 防流动性差的假胜利; 止损/时间止损看参考价 cur_price — 防盘口瞬时蒸发误触发。\nbest_bid 拉不到时自动退回参考价, 不会卡死。',
        bar=NAVY, bg=CHIP_BG)

# ---------- 7. 重评触发 ----------
s = new_slide(prs, '05 · 自动重评', '自动重评 ① — 什么时候触发')
callout(s, ML, Inches(1.2), CW, Inches(0.82),
        '一句话: 亏到接近止损线时不盲卖, 先让 AI 联网重查一遍再决定',
        '触发点设在止损线前 5 个点, 给 AI 留出反应时间。',
        bar=ACCENT, bg=BLUE_BG)
table(s, ML, Inches(2.2), CW, [3.0, 3.4, 5.83], [
    ['仓位类型', '亏到多少触发重评', '备注'],
    ['收敛型 convergent', '−15%', '止损线在 −20% (移动止损确认后也先交重评)'],
    ['混合型 hybrid', '−30%', '止损线 = 从最高价回撤 35% (v7.4.4 移动止损)'],
    ['事件型 event_driven', '−30%', '本身有 −60% 很松的止损线 (v7.4.3); 重评说卖还要过护栏 (P8)'],
    ['未分类 (少见)', '−30%', '默认当混合型, 止损线=回撤35% (不再有 −25% 老仓档)'],
], fs=11.8, row_h=0.44)
bullets(s, ML, Inches(4.85), CW, [
    ('防重复烧钱三道闸: ', '① 该仓已有重评在跑 → 锁住 ② 每仓 6 小时冷却 ③ 冷却结束不放炮 — 记下当时亏损当基线, 从基线再多亏 ≥10 个点才第二次触发'),
    ('手动 🤖 API 重评按钮: ', '绕过所有节流立刻重评, 但结果永远等你确认'),
    ('哪里看: ', '主页「📜 重评历史 & 冷却状态」折叠栏 — 每仓 ❄️ 冷却倒计时+历史决策; 副屏持仓行也有 ❄️ 徽章'),
    ('活例子: ', 'MBS 仓 (事件型) 跌破 −30% 时已自动触发过一轮重评'),
], size=12.2, gap=8)

# ---------- 8. 重评决策与执行 ----------
s = new_slide(prs, '06 · 自动重评', '自动重评 ② — AI 怎么决定、谁来执行')
txt(s, ML, Inches(1.16), CW, Inches(0.3), [
    {'runs': [('AI 联网搜完新闻后必须四选一:', {'size': 12.5, 'bold': True, 'color': NAVY})]}])
flow(s, ML, Inches(1.52), CW, Inches(0.78), [
    ('hold', '继续拿'), ('update_q', '改胜率, 继续拿'), ('exit', '卖掉'), ('cancel_autostop', '关%止损, 只留地板'),
], fill=NAVY2, head_size=12.5, sub_size=10, arrow_zone=0.22)
callout(s, ML, Inches(2.52), Inches(6.0), Inches(1.5),
        '防呆 ① 反锚定 (治"坑底卖飞")',
        '提示词里喂给 AI「大跌前 24 小时价格中枢」(排除最近 6 小时),\n不许拿坑底现价当锚。之前两笔在坑底被卖飞 ~$16,\n病根就是 AI 把砸盘后的价格当成了"合理价"。',
        bar=GREEN, bg=GREEN_BG, body_size=11)
callout(s, Inches(6.78), Inches(2.52), Inches(6.0), Inches(1.5),
        '防呆 ② 事件型卖出护栏',
        'AI 说卖也要过闸: 必须「论点被真新闻推翻」或「AI 新胜率\n比现价低 ≥8 个点」, 否则自动降级成"改胜率继续拿"。\n数据缺失时默认不卖 — 事件型常先砸后弹, 绝不在坑底盲砍。',
        bar=GREEN, bg=GREEN_BG, body_size=11)
table(s, ML, Inches(4.28), CW, [1.9, 7.0, 3.33], [
    ['你的状态', '重评结果怎么处理', '会不会自动卖'],
    ['在线', '结果挂主页红闪卡等你 ✅ 确认; 闪 2 分钟没人理 → 自动调 API, 结果仍等确认', '绝不自动卖'],
    ['离线', 'AI 决策直接自动执行 — exit 会真的卖', '会 ⚠️ 此路径未实测'],
    ['手动 🤖', '不分在线离线, 一律挂卡等确认 (出门快速一键重评用)', '绝不自动卖'],
], fs=11.3, row_h=0.52, cell_over={(2, 2): {'color': RED, 'bold': True}})
bullets(s, ML, Inches(6.42), CW, [
    ('在线/离线怎么判定: ', '页面 10 分钟没操作 → 弹「还在吗」→ 60 秒没答 → 转离线; 之后任何操作自动回在线 (手动点的离线除外)'),
], size=11.4, gap=5)

# ---------- 9. 模型 ----------
s = new_slide(prs, '07 · 自动重评', '自动重评 ③ — 用哪个 AI')
table(s, ML, Inches(1.3), CW, [2.5, 3.2, 6.53], [
    ['角色', '模型', '说明'],
    ['权威 (说了算)', 'Claude Opus 4.8', '联网搜索 + 深度思考 (省钱档 medium); 决策一律以它为准'],
    ['备用 (兜底)', '智谱 GLM-5.2', 'Claude 挂了才顶上; 思考拉满 + 专业搜索, 解析带正则兜底'],
], fs=12, row_h=0.52)
bullets(s, ML, Inches(3.2), CW, [
    ('双跑对比 (默认开): ', '每次重评两家并行都跑, Claude 说了算; GLM 结果只进 /api_reeval 对比页 — 看动作一致率 / 胜率差距 / 谁更激进 / 谁失败得多。⚠️ 真实双跑至今还没发生过'),
    ('安全底线: ', '任一模型给的决策不合法 → 直接丢弃, 绝不进执行; 两家全失败 → 不动仓位'),
    ('紧急暂停: ', '主页开关一键停付费 API; 暂停期间也不会退回盲卖 — 仓位冻结在「等重评」等你处理'),
    ('提示词同源: ', 'API 重评用的就是手动复制那份提示词 (改 prompts.py 两边同时生效), 只是各自附加决策指令'),
], size=12.4, gap=10)
callout(s, ML, Inches(5.75), CW, Inches(1.05),
        '钱花在哪',
        '只有自动重评和手动 🤖 按钮会调付费 API (Claude / 智谱)。找仓分析 = 你手动复制提示词贴到 Claude.ai 网页, 免费;\n测试仓全程零 API。在线时大跌也先出手动卡, 不急着烧 API (2 分钟没人理才升级)。',
        bar=GOLD, bg=AMBER_BG)

# ---------- 10. 金额 ----------
s = new_slide(prs, '08 · 仓位管理', '新仓金额 — 一笔下多少钱')
bullets(s, ML, Inches(1.24), CW, [
    {'m': '①', 'mc': GOLD, 'b': '先看有没有便宜可占: ', 't': 'edge = AI 胜率 q − 市场价 p; edge ≤ 0 → 直接 $0 不下'},
    {'m': '②', 'mc': GOLD, 'b': '凯利公式打 2.5 折: ', 't': '金额 = 总资产 × [edge ÷ (1−p)] × 25%'},
    {'m': '③', 'mc': GOLD, 'b': '再打两个折扣: ', 't': '结算太远 (>21 天) 按平方根比例减、最低 4 折; 冷门低价盘 (p<$0.15) 最多打对折'},
    {'m': '④', 'mc': GOLD, 'b': '三道帽子: ', 't': '同主题簇 ≤ 总资产 20% · 月回撤预算 $30 还剩多少 ÷ 该档预期跌幅 (收敛 20%/混合 35%/事件 70%) · 帽完不足 $1 → 放弃'},
    {'m': '⑤', 'mc': GOLD, 'b': '硬边界: ', 't': '最终金额压回 $1 – $15 之间'},
], size=13, gap=14)
callout(s, ML, Inches(4.25), CW, Inches(1.02),
        '算一笔现成的 (数字照真实公式)',
        '总资产 $93 · AI 胜率 60% · 市场价 $0.50 · 24 天后结算 → edge 10 点 → 凯利 $93 × 0.2 × 25% ≈ $4.7\n→ 天数打 94 折 ≈ $4.4 → 各道帽子都没碰到 → 建议下 $4.40',
        bar=ACCENT, bg=BLUE_BG)
bullets(s, ML, Inches(5.55), CW, [
    ('信心 high/medium/low 只做记录不进公式; ', 'AI 的胜率 q 在发现阶段已向市场价折过一半, 公式不再重复打折'),
    ('每次算的建议都落库 (sizing_log): ', '和你实际下的金额对比, 攒 2–4 周数据回头调参'),
    ('入口: ', '主页 JSON 通道「💵 填入计算器」自动带入这套公式; 测试仓录入用同一公式 (但不占簇额度)'),
], size=12.2, gap=8)

# ---------- 11. 找新仓 ----------
s = new_slide(prs, '09 · 找仓', '找新仓 — 从扫描到录入')
table(s, ML, Inches(1.26), CW, [2.0, 2.6, 2.5, 2.4, 2.73], [
    ['扫描档位', '价格区间', '7 天成交', '距结算', '滑点上限 ($5 单)'],
    ['标准', '$0.08 – 0.92', '≥ $1,500', '7 – 21 天', '≤ 1.5 点'],
    ['中范围', '$0.05 – 0.95', '≥ $500', '5 – 35 天', '≤ 3 点'],
    ['大范围', '$0.03 – 0.97', '≥ $200', '3 – 60 天', '≤ 5 点'],
], fs=11.8, row_h=0.42, col_align=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.CENTER])
bullets(s, ML, Inches(3.28), CW, [
    ('主题从哪来: ', '39 个白名单主题 + /tags 动态热门榜 (跟成交量热度走: 俄乌→伊朗→世界杯); 结算规则含糊的市场按关键词直接排除'),
], size=12.2, gap=6)
txt(s, ML, Inches(3.95), CW, Inches(0.3), [
    {'runs': [('固定流程 (全程约 5 分钟):', {'size': 12.5, 'bold': True, 'color': NAVY})]}])
flow(s, ML, Inches(4.3), CW, Inches(0.95), [
    ('① 一键全扫', '5 主题并发\n约 20 秒'),
    ('② 拼提示词', '报告+持仓+簇字典\n自动拼好'),
    ('③ 贴给 Claude', 'Claude.ai 网页\n免费分析'),
    ('④ 回 JSON 推荐', '方向/q/信心\n止损档/簇'),
    ('⑤ 算金额', '💵 一键带入\n计算器'),
    ('⑥ 下单+录入', '手动买 → 📌\n一键录全'),
], head_size=11, sub_size=8.8, arrow_zone=0.26)
bullets(s, ML, Inches(5.55), CW, [
    ('JSON 契约: ', 'Claude 每条推荐固定给全 方向/胜率 q/信心/止损档/主题簇/理由 — 录入时一键写全, bot 靠这些盯盘'),
    ('拿不准的推荐: ', '点 🧪 丢进测试仓白嫖验证 (P14), 不占真钱额度'),
    ('草稿不怕刷新: ', '粘贴的推荐存在浏览器里, 去下单回来接着录; 点「🗑 清理」或录入成功才消'),
], size=12.2, gap=8)

# ---------- 12. 入场门槛 (edge, prompts.py §推荐门槛 v7.2) ----------
s = new_slide(prs, '10 · 找仓', '入场门槛 — 多大的"便宜"才值得下')
callout(s, ML, Inches(1.14), CW, Inches(1.02),
        '怎么算: 便宜 (edge) 必须大过门槛, Claude 才允许推荐',
        'edge = 校准后胜率 − 下注方向现价; 校准 = 先把 AI 跟市场的分歧打五折 (市场价 + 0.5 × 分歧)。\n最终门槛 = max( 5pp 硬地板, ① 基础门槛 + ② 价位叠加 ); 5pp ≈ 滑点+双边手续费, 再低净利被执行成本吃光。够不着就"无推荐", 不凑数。',
        bar=ACCENT, bg=BLUE_BG)
txt(s, ML, Inches(2.36), Inches(5.5), Inches(0.28), [
    {'runs': [('① 基础门槛 — 按扫描范围', {'size': 12, 'bold': True, 'color': NAVY})]}])
table(s, ML, Inches(2.66), Inches(5.4), [1.5, 1.3, 2.6], [
    ['扫描档位', '基础门槛', '为什么'],
    ['标准', '6pp', '盘面最干净, 门槛最低'],
    ['中范围', '8pp', '过滤更宽 → 要更多安全垫'],
    ['大范围', '10pp', '盘子薄、成本高 → 最严'],
], fs=11.5, row_h=0.44, header_h=0.34,
  col_align=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT])
txt(s, Inches(6.35), Inches(2.36), Inches(6.2), Inches(0.28), [
    {'runs': [('② 价位叠加 — 按"你下注方向"那侧的现价', {'size': 12, 'bold': True, 'color': NAVY})]}])
table(s, Inches(6.35), Inches(2.66), Inches(6.43), [2.2, 1.4, 2.83], [
    ['下注方向现价', '调整', '为什么'],
    ['≥ 65¢ (强 favorite)', '−3pp', '方向高度可靠, 小便宜也值'],
    ['50 – 65¢', '−1pp', '略松'],
    ['35 – 50¢', '±0', '基准区'],
    ['< 35¢ (longshot)', '+3pp 更严', '冷门系统性被高估'],
], fs=11.5, row_h=0.44, header_h=0.34,
  col_align=[PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.LEFT])
bullets(s, ML, Inches(4.98), CW, [
    ('实证依据 (本 bot 已结算样本): ', '下注方向买入价 ≥70¢ 的 12 笔全对, <30¢ 的只有 15% 兑现 — 冷门被系统性高估 (favorite-longshot bias), 所以越冷门要求越苛刻'),
    ('别跟持仓状态灯搞混: ', '这套门槛只管挑新仓 (写在 DISCOVERY 提示词 §推荐门槛, Claude 挑仓时执行); 已持仓的 +2/−3 点状态灯 (P6) 是另一套'),
], size=12.2, gap=8)
callout(s, ML, Inches(6.02), CW, Inches(0.92),
        '三个现成例子 (最终门槛 = 基础 + 叠加, 但不低于 5pp)',
        '中范围买 NO @ 72¢ → 8 − 3 = 5pp 就够 · 中范围买 YES @ 25¢ → 8 + 3 = 11pp 才行 · 大范围买 @ 58¢ → 10 − 1 = 9pp。\n一句话: 价越高越松 (最松贴到 5pp 地板), 越冷门越严。',
        bar=GREEN, bg=GREEN_BG)

# ---------- 13. 钱的总规矩 ----------
s = new_slide(prs, '11 · 资金风控', '钱的总规矩 — 风控红线一页')
chips(s, ML, Inches(1.24), CW, [
    ('$1 – 15', '单仓上下限'),
    ('≤ 20%', '单主题簇 ≈ $18.7'),
    ('$30 / 月', '组合回撤预算'),
    ('$0.05', '绝对地板价'),
], h=Inches(0.95))
bullets(s, ML, Inches(2.5), CW, [
    ('买入 0 自动: ', 'bot 没有自动买入代码路径, 每一笔买单都出自你手'),
    ('自动卖出只有 3 个口: ', '① 止盈到线 ② 跌破 $0.05 地板 ③ 你离线时重评结论=卖 — 每笔卖出都落账进已平仓表'),
    ('在线/离线是动钱开关: ', '在线 = 只建议+等确认; 离线 = 自动执行 (⚠️ 未实测); 手动 🤖 按钮永远等确认'),
    ('测试仓物理隔离: ', '/paper 代码碰不到下单函数和付费 API; 改这块代码前必须先跑审计检查'),
    ('资产口径统一按现价: ', '总资产 = 现金 + Σ(现价×份数); 簇暴露/回撤预算也按现价算, 不按成本'),
    ('钱包与访问: ', '代理钱包 GNOSIS_SAFE 签名, .env 配置冻结不改; 面板仅内网 (Tailscale)+密码, 错 5 次锁 30 分钟'),
    ('数据保底: ', '每 30 分钟自动 git 备份+资产快照导出; 已平仓 80 笔 07-06 刚与 Polymarket 全量对齐 (累计 −$3.51)'),
], size=13, gap=17)

# ---------- 14. 测试仓 ----------
s = new_slide(prs, '12 · 测试仓', '测试仓 /paper — 不花钱的模拟盘')
bullets(s, ML, Inches(1.28), CW, [
    ('干嘛用: ', '拿不准/看着离谱的 AI 推荐先丢进来, 一分钱不花, 跑和真仓完全同一套盯盘算法, 验证 AI 到底准不准'),
    ('怎么进: ', '主页 JSON 通道每条推荐带「🧪 加入测试仓」一键录 (金额用真仓同款公式自动算, 算不出兜底 $10); /paper 页也能手动填'),
    ('它做什么: ', '每 30 秒同真仓一起盯 → 算法一"卖"(命中卖出条件)就移到「往期测试仓」→ 继续盯到结算给最终对错。像真仓一样有生命周期 (进行中 / 往期)'),
    ('重评也免费: ', '每仓一键复制重评提示词 → 自己贴到 Claude.ai 网页 (免费) → 把新胜率存回来; 全程零 API 零钱'),
    ('特例: ', '金额不受同簇真仓挤压 — 簇上限是真钱风控, 测试仓不占额度 (月回撤预算照常算)'),
    ('往期 + 统计 (v7.4.0): ', '往期区每条: 预测准不准 + 模拟盈亏 + 📈最高点 (本可赚多少, 哪怕后来亏到底); 顶部统计: 模拟总盈亏 / 赚钱率 / 结算对率 — 一眼看这套 AI 准不准'),
], size=13, gap=17)
callout(s, ML, Inches(5.35), CW, Inches(1.15),
        '🔒 铁律 (2026-06-22 定, 代码层隔离, 不是口头约定)',
        '测试仓永不碰钱: 不下单、不卖、不调付费 API, 只允许只读拉行情。\n真仓的自动重评天生够不到它 (独立表+独立代码路径)。',
        bar=RED, bg=RED_BG)

# ---------- 15. 参数速查 ----------
s = new_slide(prs, '13 · 参数速查', '常改参数速查表 (改完重启生效)')


def _p(zh, code):
    return {'lines': [(zh, {'size': 10.3, 'bold': True, 'color': INK}),
                      (code, {'size': 7.6, 'color': GRAY, 'font': 'Menlo'})]}


txt(s, ML, Inches(1.12), Inches(6), Inches(0.28), [
    {'runs': [('出场规则 — modules/monitor.py 顶部常量', {'size': 11.5, 'bold': True, 'color': NAVY})]}])
table(s, ML, Inches(1.42), Inches(5.95), [4.35, 1.6], [
    ['参数', '现值'],
    [_p('通用止盈 (全卖)', 'TAKE_PROFIT_PRICE'), '$0.90'],
    [_p('事件型止盈 (卖一半)', 'TAKE_PROFIT_PRICE_EVENT_DRIVEN'), '$0.92'],
    [_p('收敛临近止盈 / 天数', 'TAKE_PROFIT_PRICE_CONVERGENT_NEAR / _DAYS'), '$0.88 / 3 天'],
    [_p('翻倍止盈 (事件型除外)', 'TAKE_PROFIT_PNL_PCT'), '+100%'],
    [_p('移动止损 收敛/临近/混合/确认', 'TRAILING_STOP_PCT_CONVERGENT/_NEAR/_HYBRID/ROUNDS'), '20% / 12% / 35% / 6'],
    [_p('三档止损 (收敛/混合/事件)', 'STOP_LOSS_PCT_BY_TIER'), '20% / 35% / 60%'],
    [_p('未分类默认 / 地板价', '→当hybrid / EVENT_DRIVEN_FLOOR'), '当混合35% / $0.05'],
    [_p('时间止损 天数 / 漂移', 'TIME_STOP_DAYS / _DRIFT_PP'), '2 天 / 5 点'],
    [_p('状态灯 持有 / 偏空', 'HOLD_MIN_EDGE_PP / SOFT_NEGATIVE'), '+2 / −3 点'],
    [_p('盯盘心跳', 'CHECK_INTERVAL'), '30 秒'],
], fs=10.3, row_h=0.46, header_h=0.32, col_align=[PP_ALIGN.LEFT, PP_ALIGN.CENTER])
txt(s, Inches(6.85), Inches(1.12), Inches(6), Inches(0.28), [
    {'runs': [('重评 & 金额 — .env 环境变量 (不用改代码)', {'size': 11.5, 'bold': True, 'color': NAVY})]}])
table(s, Inches(6.85), Inches(1.42), Inches(5.93), [4.33, 1.6], [
    ['参数', '现值'],
    [_p('重评冷却', 'AUTO_REEVAL_COOLDOWN_H'), '6 小时'],
    [_p('冷却后再触发需再亏', 'AUTO_REEVAL_RETRIGGER_DROP'), '10 点'],
    [_p('在线红闪升级等待', 'AUTO_REEVAL_MANUAL_ESCALATE_MIN'), '2 分钟'],
    [_p('事件型卖出护栏', 'AUTO_REEVAL_EXIT_GUARD_EDGE'), '8 点'],
    [_p('反锚定中枢窗口 / 排除', 'AUTO_REEVAL_CENTER_WINDOW_H / SKIP_H'), '24h / 6h'],
    [_p('权威模型 / 双跑对比', 'AUTO_REEVAL_PRIMARY / _DUAL'), 'claude / 开'],
    [_p('凯利分数', 'SIZING_KELLY_FRACTION'), '0.25'],
    [_p('月回撤预算', 'SIZING_MONTHLY_DD_BUDGET'), '$30'],
    [_p('单簇上限', 'SIZING_CLUSTER_CAP_PCT'), '20%'],
    [_p('单仓上限 / 下限', 'SIZING_MAX / MIN_SINGLE_POS'), '$15 / $1'],
], fs=10.3, row_h=0.46, header_h=0.32, col_align=[PP_ALIGN.LEFT, PP_ALIGN.CENTER])
txt(s, ML, Inches(6.58), CW, Inches(0.42), [
    {'runs': [('※ 改完必须重启 bot 才生效 (CLAUDE.md 有整段重启命令), 重启后看 bot.log 无报错才算成功; 版本号只改 modules/version.py 一处, 全站自动同步。',
               {'size': 9.3, 'color': GRAY})], 'space_after': 2},
    {'runs': [('※ 入场 edge 门槛 (基础 6/8/10pp ± 价位叠加, 地板 5pp, 见 P12) 不在上表 — 在 modules/prompts.py §推荐门槛 (DISCOVERY 提示词) 里改, 同样重启生效。',
               {'size': 9.3, 'color': GRAY})]}])

# ---------- 16. 风险 ----------
s = new_slide(prs, '14 · 风险备忘', '已知风险与注意事项')
callout(s, ML, Inches(1.22), CW, Inches(1.0),
        '⚠️ 出场阈值是"方向值", 不是精调值',
        '0.92 / 0.88 / 20% / 12% / 8 点 / 6 拍全部来自 51 笔小样本回测。每次重评都在落校准数据 (大跌前中枢+价格曲线),\n攒够了再统一精调 — 在那之前别对单笔结果过度反应。',
        bar=AMBER, bg=AMBER_BG)
callout(s, ML, Inches(2.4), CW, Inches(1.18),
        '🔴 6·26 文件回退事故 (已恢复, 根因未明)',
        'monitor + auto_reeval 曾被莫名整体回退到老版, v7.0 出场策略停摆 10 天 (06-26 → 07-06), 期间实际跑的是\n老版"0.90 一刀切全卖 + 盲止损"。07-06 已从 git 完整找回并提交锁定。若再发现行为像老版, 先查这两个文件。',
        bar=RED, bg=RED_BG)
bullets(s, ML, Inches(3.85), CW, [
    ('两条路径从未实测: ', '① 离线自动卖出 (代码在, 没真跑过) ② Claude+GLM 双跑对比 — 第一次出现时盯着 bot.log 看'),
    ('数据口径: ', '已平仓表 07-06 按 Polymarket 真实成交重建 (80 笔, 累计 −$3.51); 已加守卫防"部分卖出记成假平仓"'),
    ('目前成绩单 (样本小, 仅参考): ', '止盈规则 11 笔全部赚钱 +$25 🟢 · 重评清仓 14 笔合计 −$8.6 🔴 · 历史"卖飞"共少赚 $42.48 — v7.0 出场重设计就是治这个'),
    ('下一步: ', '用测试仓+校准数据验证阈值 → 实测离线卖出与双跑 → 按真实数据调参'),
], size=13, gap=18)

prs.save(OUT)
print(f'OK -> {OUT}')
print('slides:', len(prs.slides._sldIdLst))
